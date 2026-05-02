#!/usr/bin/env python3
"""
基于附件3-2模板PPT和开发与应用报告，生成演示PPT
结构：封面 → 案例概述 → 实现功能总览(5项) → 5项功能分页(记忆系统3页) → 应用情况 → 创新与展望 → 感谢
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import os

ASSETS = '/workspace/projects/assets'
OUTPUT = os.path.join(ASSETS, '演示PPT.pptx')

# ========== 工具函数 ==========

def add_textbox(slide, left, top, width, height, items):
    """添加文本框，items = [(text, bold, size, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text = item[0]
        bold = item[1] if len(item) > 1 else False
        size = item[2] if len(item) > 2 else Pt(17)
        color = item[3] if len(item) > 3 else RGBColor(0x33, 0x33, 0x33)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = '楷体'
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(3)
    return txBox

def add_screenshot_placeholder(slide, left, top, width, height, label):
    """添加截图占位区域（虚线框）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    line = shape.line
    line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    line.width = Pt(1.5)
    line.dash_style = 2  # Dash
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = '楷体'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_image_centered(slide, img_path, top, max_width=Inches(9.5), max_height=Inches(4.5)):
    img = Image.open(img_path)
    img_w, img_h = img.size
    ratio = min(max_width / Inches(img_w / 96), max_height / Inches(img_h / 96))
    final_w = int(Inches(img_w / 96) * ratio)
    final_h = int(Inches(img_h / 96) * ratio)
    left = (12192000 - final_w) // 2
    slide.shapes.add_picture(img_path, left, top, final_w, final_h)

def add_title_bar(slide, title_text):
    """添加蓝色标题栏"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 12192000, Emu(850000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x2C, 0x48, 0x9B)
    bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = '黑体'
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Emu(500000)
    tf.margin_top = Emu(150000)

# ========== 生成示意图 ==========

def create_voice_diagram():
    W, H = 1100, 380
    img = Image.new('RGB', (W, H), '#FAFAFA')
    draw = ImageDraw.Draw(img)
    try:
        ft = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 18)
        fl = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 14)
        fs = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 12)
    except:
        ft = fl = fs = ImageFont.load_default()

    draw.text((W//2, 15), "语音交互技术流程", fill='#2C489B', font=ft, anchor='mt')
    nodes = [('麦克风\n采集PCM', '#2C489B'), ('VAD\n语音检测', '#4A7C59'),
             ('ASR\n语音转文字', '#7B4FA0'), ('LLM+RAG\n智能回答', '#C49A6C'),
             ('TTS\n文字转语音', '#2C489B'), ('音频波纹\n播放语音', '#4A7C59')]
    for i, (label, color) in enumerate(nodes):
        x = 60 + i * 180; y = 50
        draw.rounded_rectangle([x, y, x+140, y+80], radius=10, fill=color)
        for j, line in enumerate(label.split('\n')):
            draw.text((x+70, y+28+j*20), line, fill='white', font=fs, anchor='mm')
        if i < 5:
            ax = x + 150
            draw.line([(ax, y+40), (ax+25, y+40)], fill='#666', width=2)
            draw.polygon([(ax+22, y+35), (ax+30, y+40), (ax+22, y+45)], fill='#666')

    y = 170
    draw.text((W//2, y), "关键创新", fill='#8B0000', font=ft, anchor='mt')
    innovations = [('回声防护', 'TTS播放时静音麦克风\n(gain=0)，杜绝回声'),
                   ('TTS预解码', '预解码AudioBuffer\n队列，无缝播放'),
                   ('短句合并', '<8字短句暂不发送\n减少无效TTS请求'),
                   ('手动打断', '脉冲停止按钮\n即时中断切换状态')]
    for i, (title, desc) in enumerate(innovations):
        x = 60 + i * 260; yb = y + 30
        draw.rounded_rectangle([x, yb, x+240, yb+130], radius=8, fill='white', outline='#C49A6C', width=1)
        draw.rounded_rectangle([x, yb, x+240, yb+35], radius=8, fill='#2C489B')
        draw.rectangle([x, yb+25, x+240, yb+35], fill='#2C489B')
        draw.text((x+120, yb+17), title, fill='white', font=fl, anchor='mm')
        for j, line in enumerate(desc.split('\n')):
            draw.text((x+120, yb+65+j*22), line, fill='#555', font=fs, anchor='mm')

    path = os.path.join(ASSETS, 'diagram-voice.png')
    img.save(path, quality=95)
    return path

def create_memory_overview_diagram():
    """学生记忆系统总览图：四维数据 + 三阶段工作流"""
    W, H = 1100, 520
    img = Image.new('RGB', (W, H), '#FAFAFA')
    draw = ImageDraw.Draw(img)
    try:
        ft = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 16)
        fl = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 14)
        fs = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 12)
        fss = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 11)
    except:
        ft = fl = fs = fss = ImageFont.load_default()

    # 四大数据维度
    draw.text((W//2, 12), "四大数据维度", fill='#2C489B', font=ft, anchor='mt')
    dimensions = [
        ('学生画像', '学习风格\n兴趣偏好', '#2C489B'),
        ('知识掌握', '三档分类\n量化评分', '#C49A6C'),
        ('对话记录', '历史摘要\n困惑点', '#4A7C59'),
        ('教学策略', '有效方法\n学习建议', '#7B4FA0'),
    ]
    for i, (title, desc, color) in enumerate(dimensions):
        x = 30 + i * 270
        draw.rounded_rectangle([x, 38, x+245, 155], radius=8, fill='white', outline=color, width=2)
        draw.rounded_rectangle([x, 38, x+245, 68], radius=8, fill=color)
        draw.rectangle([x, 58, x+245, 68], fill=color)
        draw.text((x+122, 53), title, fill='white', font=fl, anchor='mm')
        for j, line in enumerate(desc.split('\n')):
            draw.text((x+122, 90+j*22), line, fill='#555', font=fs, anchor='mm')

    # 箭头：数据维度 → 工作流
    draw.line([(W//2, 165), (W//2, 185)], fill='#666', width=2)
    draw.polygon([(W//2-6, 180), (W//2, 192), (W//2+6, 180)], fill='#666')

    # 三阶段工作流
    draw.text((W//2, 200), "三阶段工作流", fill='#2C489B', font=ft, anchor='mt')
    stages = [
        ('对话前：记忆检索', ['加载学生画像', '获取知识掌握状态', '检索教学策略记忆', '构建记忆摘要', '注入系统提示词'], '#2C489B'),
        ('对话中：个性化回复', ['已掌握不重复讲解', '薄弱点重点讲解', '适配学习风格', '出验证测试题', '价值观领域约束'], '#C49A6C'),
        ('对话后：记忆更新', ['LLM自动分析对话', '更新知识掌握程度', '提取困惑与薄弱点', '生成学习建议', '完善学生画像'], '#4A7C59'),
    ]
    for i, (title, items, color) in enumerate(stages):
        x = 30 + i * 370
        draw.rounded_rectangle([x, 225, x+340, 505], radius=10, fill='white', outline=color, width=2)
        draw.rounded_rectangle([x, 225, x+340, 260], radius=10, fill=color)
        draw.rectangle([x, 248, x+340, 260], fill=color)
        draw.text((x+170, 242), title, fill='white', font=fl, anchor='mm')
        for j, item in enumerate(items):
            y = 275 + j * 44
            draw.rounded_rectangle([x+15, y, x+325, y+36], radius=6, fill='#F5F5F5', outline=color, width=1)
            draw.text((x+170, y+18), item, fill='#333', font=fss, anchor='mm')
        if i < 2:
            ax = x + 350
            draw.line([(ax, 365), (ax+12, 365)], fill=color, width=3)
            draw.polygon([(ax+8, 359), (ax+16, 365), (ax+8, 371)], fill=color)

    path = os.path.join(ASSETS, 'diagram-memory-overview.png')
    img.save(path, quality=95)
    return path

def create_mastery_diagram():
    """知识掌握三档评价图"""
    W, H = 1100, 300
    img = Image.new('RGB', (W, H), '#FAFAFA')
    draw = ImageDraw.Draw(img)
    try:
        fl = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 16)
        fs = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 13)
    except:
        fl = fs = ImageFont.load_default()
    levels = [
        ('已掌握 (>=0.6)', '#4A7C59', '不重复讲解基础\n可进阶讨论', '通过验证测试\nmastery=0.8'),
        ('学习中 (0.3-0.5)', '#C49A6C', '有了解未完全掌握\n需巩固和补充', '说懂了但未测试\n+0.1'),
        ('薄弱 (<0.3)', '#8B0000', '理解困难\n需重点讲解', '困惑或答错\n-0.1'),
    ]
    for i, (title, color, strategy, update) in enumerate(levels):
        x = 20 + i * 360
        draw.rounded_rectangle([x, 15, x+340, 285], radius=10, fill='white', outline=color, width=2)
        draw.rounded_rectangle([x, 15, x+340, 55], radius=10, fill=color)
        draw.rectangle([x, 42, x+340, 55], fill=color)
        draw.text((x+170, 35), title, fill='white', font=fl, anchor='mm')
        draw.text((x+170, 72), "教学策略", fill='#666', font=fs, anchor='mt')
        for j, line in enumerate(strategy.split('\n')):
            draw.text((x+170, 95+j*22), line, fill='#333', font=fs, anchor='mt')
        draw.line([(x+20, 150), (x+320, 150)], fill='#DDD', width=1)
        draw.text((x+170, 160), "评价规则", fill='#666', font=fs, anchor='mt')
        for j, line in enumerate(update.split('\n')):
            draw.text((x+170, 183+j*22), line, fill='#333', font=fs, anchor='mt')
    path = os.path.join(ASSETS, 'diagram-mastery.png')
    img.save(path, quality=95)
    return path

def create_diagnosis_diagram():
    """学习诊断与建议示意图"""
    W, H = 1100, 280
    img = Image.new('RGB', (W, H), '#FAFAFA')
    draw = ImageDraw.Draw(img)
    try:
        fl = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 15)
        fs = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 12)
        fss = ImageFont.truetype('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc', 11)
    except:
        fl = fs = fss = ImageFont.load_default()

    # 发现的问题
    x1 = 30
    draw.rounded_rectangle([x1, 15, x1+500, 265], radius=10, fill='white', outline='#8B0000', width=2)
    draw.rounded_rectangle([x1, 15, x1+500, 50], radius=10, fill='#8B0000')
    draw.rectangle([x1, 38, x1+500, 50], fill='#8B0000')
    draw.text((x1+250, 32), "发现的问题（对话后LLM自动提取）", fill='white', font=fl, anchor='mm')
    problems = [
        ("困惑点", "对话中识别的难点和误区，如\"不理解加速度概念\""),
        ("薄弱项", "连续答错或理解困难的知识点，mastery_level自动降低"),
        ("错误概念", "学生持有的错误理解，如\"混淆速度和加速度\""),
    ]
    for j, (title, desc) in enumerate(problems):
        y = 65 + j * 65
        draw.rounded_rectangle([x1+15, y, x1+485, y+55], radius=6, fill='#FFF5F5', outline='#8B0000', width=1)
        draw.text((x1+35, y+12), title, fill='#8B0000', font=fs, anchor='lm')
        draw.text((x1+35, y+35), desc, fill='#555', font=fss, anchor='lm')

    # 学习建议
    x2 = 570
    draw.rounded_rectangle([x2, 15, x2+500, 265], radius=10, fill='white', outline='#4A7C59', width=2)
    draw.rounded_rectangle([x2, 15, x2+500, 50], radius=10, fill='#4A7C59')
    draw.rectangle([x2, 38, x2+500, 50], fill='#4A7C59')
    draw.text((x2+250, 32), "学习建议（自动生成并注入AI上下文）", fill='white', font=fl, anchor='mm')
    suggestions = [
        ("下次方向", "基于薄弱点推荐后续学习内容，避免无目的重复"),
        ("有效方法", "记住对学生最有效的讲解方式（如实验类比、图示法）"),
        ("突破巩固", "已取得进展的知识点需要继续强化，防止回退"),
    ]
    for j, (title, desc) in enumerate(suggestions):
        y = 65 + j * 65
        draw.rounded_rectangle([x2+15, y, x2+485, y+55], radius=6, fill='#F0FFF0', outline='#4A7C59', width=1)
        draw.text((x2+35, y+12), title, fill='#4A7C59', font=fs, anchor='lm')
        draw.text((x2+35, y+35), desc, fill='#555', font=fss, anchor='lm')

    # 中间箭头
    draw.text((545, 140), "→", fill='#2C489B', font=fl, anchor='mm')

    path = os.path.join(ASSETS, 'diagram-diagnosis.png')
    img.save(path, quality=95)
    return path


# ========== 主逻辑 ==========

def delete_all_slides(prs):
    """删除演示文稿中的所有幻灯片（通过XML操作）"""
    from lxml import etree
    
    # 获取presentation XML元素
    pres_element = prs._element
    # 找到p:sldIdLst元素
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    sldIdLst = pres_element.find('.//p:sldIdLst', nsmap)
    if sldIdLst is None:
        print("  WARNING: sldIdLst not found")
        return
    
    sldIds = list(sldIdLst)
    for sldId in sldIds:
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
    print(f"  Deleted {len(sldIds)} slides")

def main():
    print("Generating diagrams...")
    voice_path = create_voice_diagram()
    memory_overview_path = create_memory_overview_diagram()
    mastery_path = create_mastery_diagram()
    diagnosis_path = create_diagnosis_diagram()

    # 基于模板创建，保留模板的元数据和主题
    template_path = os.path.join(ASSETS, 'template.pptx')
    prs = Presentation(template_path)
    print(f"  Loaded template: {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")
    
    # 删除模板中的所有幻灯片
    delete_all_slides(prs)
    print(f"  Cleared template slides")

    # 使用模板的布局（布局0=标题幻灯片, 布局6=空白, 其他按模板定义）
    # 找到空白布局
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == '空白' or layout.name == 'Blank':
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    # 颜色常量
    BLUE = RGBColor(0x2C, 0x48, 0x9B)
    DARK = RGBColor(0x33, 0x33, 0x33)
    RED = RGBColor(0x8B, 0x00, 0x00)
    BROWN = RGBColor(0xC4, 0x9A, 0x6C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ========================================
    # Slide 1: 封面
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 12192000, 6858000)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x2C, 0x48, 0x9B)
    bg.line.fill.background()
    add_textbox(slide, Emu(800000), Emu(1800000), Emu(10600000), Emu(2000000),
                [('开放智慧助教', True, Pt(44), WHITE),
                 ('基于教师知识库的智能语音交互教学系统', False, Pt(24), RGBColor(0xDD, 0xDD, 0xFF))])
    add_textbox(slide, Emu(800000), Emu(4800000), Emu(10600000), Emu(1500000),
                [('XX大学继续教育学院', False, Pt(20), RGBColor(0xCC, 0xCC, 0xEE)),
                 ('团队成员：李明  王芳  陈刚', False, Pt(18), RGBColor(0xBB, 0xBB, 0xDD))])
    print("  Slide 1: 封面")

    # ========================================
    # Slide 2: 案例概述
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '案例概述')
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(5500000),
                [('在成人高等教育领域，非全日制成人学生普遍面临工学矛盾——工作日全职上班，难以到校与教师常态化交流，课后遇到问题无人解答。同时，成人学生基础普遍薄弱，学习需求差异大，需要更细致耐心的个性化指导。传统在线答疑教师时间有限，无法持续跟踪学生认知状态，而通用AI辅导产品缺乏专业性和价值观约束，存在教育安全风险。', False, Pt(20), DARK),
                 ('', False, Pt(12), DARK),
                 ('本系统——"开放智慧助教"，借助国产大模型和智能体开发平台，让每位教师创建专属知识库和教学风格的数字助教，为成人学生提供7×24小时语音辅导，引入学生记忆系统实现"记住学生、因材施教"，并通过严格的领域约束机制确保AI助教始终聚焦专业教学，不回答与学科无关的问题，保障教育的正确价值导向。', False, Pt(20), DARK)])
    print("  Slide 2: 案例概述")

    # ========================================
    # Slide 3: 实现功能总览（5项）
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '实现功能总览')

    features_summary = [
        ('一、语音实时交互', True, Pt(22), BLUE),
        ('基于Web Audio API的实时语音问答，含回声防护、TTS预解码、短句合并、手动打断', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('二、价值观领域约束', True, Pt(22), BLUE),
        ('AI助教只回答专业相关问题，超出范围礼貌引导回擅长领域，守护教育安全', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('三、学生记忆系统', True, Pt(22), BLUE),
        ('四维追踪（画像/知识掌握/对话历史/教学策略）+ 知识掌握评价 + 学习诊断与建议', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('四、知识库管理', True, Pt(22), BLUE),
        ('教师上传文本/URL/文件，RAG语义检索增强对话', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('五、教师管理与数字人', True, Pt(22), BLUE),
        ('4步设置向导、管理面板、真人照片+音频波纹动画', False, Pt(16), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(5500000), features_summary)
    print("  Slide 3: 实现功能总览")

    # ========================================
    # Slide 4: 一、语音实时交互系统
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '一、语音实时交互系统')
    items = [
        ('基于Web Audio API构建持久化麦克风管道，实现实时语音检测闭环', False, Pt(18), DARK),
        ('麦克风 → VAD → ASR → LLM+RAG → TTS → 音频波纹动画', True, Pt(16), BROWN),
        ('', False, Pt(4), DARK),
        ('关键创新：', True, Pt(18), BLUE),
        ('• 回声防护：TTS播放时静音麦克风(gain=0)，杜绝AI语音回传被误识别', False, Pt(17), DARK),
        ('• TTS预解码：预解码AudioBuffer队列播放，消除段间延迟实现无缝语音', False, Pt(17), DARK),
        ('• 短句合并：<8字短句暂不发送TTS，减少无效请求，提升语音连贯性', False, Pt(17), DARK),
        ('• 手动打断：脉冲停止按钮，即时中断TTS恢复麦克风', False, Pt(17), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(2600000), items)
    add_image_centered(slide, voice_path, Emu(3600000), max_width=Inches(9), max_height=Inches(2.0))
    add_screenshot_placeholder(slide, Emu(1500000), Emu(5700000), Emu(9000000), Emu(900000),
                               '【请插入：语音对话界面截图（含数字人+音频波纹动画+音量指示器）】')
    print("  Slide 4: 语音实时交互系统")

    # ========================================
    # Slide 5: 二、价值观领域约束机制
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '二、价值观领域约束机制')
    items = [
        ('系统核心设计原则：AI助教只回答专业领域问题', True, Pt(20), RED),
        ('', False, Pt(6), DARK),
        ('约束机制：', True, Pt(18), BLUE),
        ('• 系统提示词严格限制：只回答与教师专业领域和知识库相关的问题', False, Pt(17), DARK),
        ('• 超出范围时礼貌说明："这个问题超出了我的专业范围，我们继续聊XX吧"', False, Pt(17), DARK),
        ('• 教师档案明确限定专业范围，AI严格按档案定义身份', False, Pt(17), DARK),
        ('• 区别于通用AI"有问必答"，坚守教育者角色，守护教学安全', False, Pt(17), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(5500000), Emu(5200000), items)
    add_screenshot_placeholder(slide, Emu(6500000), Emu(1200000), Emu(5200000), Emu(5000000),
                               '【请插入：价值观约束对话截图\n（如学生问天气被拒绝\n并引导回学科内容）】')
    print("  Slide 5: 价值观领域约束机制")

    # ========================================
    # Slide 6: 三、学生记忆系统 — 系统总览
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '三、学生记忆系统 — 总览')
    items = [
        ('核心能力："记住学生、因材施教"，按教师完全隔离', True, Pt(18), BLUE),
        ('', False, Pt(2), DARK),
        ('四大数据维度：学生画像 | 知识掌握 | 对话记录 | 教学策略', True, Pt(16), BROWN),
        ('三阶段工作流：对话前检索 → 对话中个性化 → 对话后自动更新', True, Pt(16), BROWN),
        ('两大核心功能：知识掌握评价机制 + 学习诊断与建议', True, Pt(16), RED),
    ]
    add_textbox(slide, Emu(600000), Emu(1000000), Emu(11000000), Emu(1200000), items)
    add_image_centered(slide, memory_overview_path, Emu(2100000), max_width=Inches(9.5), max_height=Inches(3.5))
    add_screenshot_placeholder(slide, Emu(1500000), Emu(5700000), Emu(9000000), Emu(900000),
                               '【请插入：教师管理面板-学生记忆总览截图\n（知识掌握三档列表+对话历史+教学策略）】')
    print("  Slide 6: 学生记忆系统-总览")

    # ========================================
    # Slide 7: 三、学生记忆系统 — 知识掌握评价机制
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '三、学生记忆系统 — 知识掌握评价')
    items = [
        ('系统在对话中主动提出验证问题，根据学生回答评价掌握程度', True, Pt(18), BLUE),
        ('', False, Pt(2), DARK),
        ('三档分类：已掌握(>=0.6) / 学习中(0.3-0.5) / 薄弱(<0.3)', True, Pt(17), BROWN),
        ('评价规则：验证通过→0.8 | 说懂未测试→+0.1 | 困惑答错→-0.1', True, Pt(17), BROWN),
        ('学生说"懂了"必须出验证测试，通过才标记已掌握', True, Pt(16), RED),
        ('知识点模糊匹配与自动合并：避免同一知识点重复记录', False, Pt(15), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(1500000), items)
    add_image_centered(slide, mastery_path, Emu(2500000), max_width=Inches(9.5), max_height=Inches(2.2))
    add_screenshot_placeholder(slide, Emu(1500000), Emu(5700000), Emu(9000000), Emu(900000),
                               '【请插入：验证测试对话截图\n（学生说懂了→AI出验证题→学生答对/答错）】')
    print("  Slide 7: 学生记忆系统-知识掌握评价")

    # ========================================
    # Slide 8: 三、学生记忆系统 — 学习诊断与建议
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '三、学生记忆系统 — 学习诊断与建议')
    items = [
        ('对话后LLM自动分析对话内容，提取问题并生成建议', True, Pt(18), BLUE),
        ('', False, Pt(2), DARK),
        ('发现的问题 → 学习建议 → 全部注入AI上下文，指导下一次对话', True, Pt(17), RED),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(800000), items)
    add_image_centered(slide, diagnosis_path, Emu(1800000), max_width=Inches(9.5), max_height=Inches(2.5))
    add_screenshot_placeholder(slide, Emu(1500000), Emu(5700000), Emu(9000000), Emu(900000),
                               '【请插入：教师管理面板-学生记忆详情截图\n（困惑点+薄弱项+学习建议列表）】')
    print("  Slide 8: 学生记忆系统-学习诊断与建议")

    # ========================================
    # Slide 9: 四、知识库管理
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '四、知识库管理')
    items = [
        ('教师专属知识库，RAG架构语义检索增强对话', True, Pt(18), BLUE),
        ('', False, Pt(4), DARK),
        ('• 文本添加：直接输入教学内容', False, Pt(17), DARK),
        ('• URL导入：从网页自动提取内容', False, Pt(17), DARK),
        ('• 文件上传：支持 .txt/.md/.csv/.json/.html/.xml/.docx/.pdf', False, Pt(17), DARK),
        ('• 语义搜索：基于向量检索浏览知识库内容', False, Pt(17), DARK),
        ('• LRU缓存：重复查询零延迟，首次对话减少200~500ms', False, Pt(17), DARK),
        ('• 教师档案：定义专业范围、教学风格、引导问题', False, Pt(17), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(5500000), Emu(5200000), items)
    add_screenshot_placeholder(slide, Emu(6500000), Emu(1200000), Emu(5200000), Emu(5000000),
                               '【请插入：知识库管理界面截图\n（添加文本/上传文件/语义搜索）】')
    print("  Slide 9: 知识库管理")

    # ========================================
    # Slide 10: 五、教师管理与数字人
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '五、教师管理与数字人')
    items = [
        ('教师4步设置向导，5分钟完成专属助教创建', True, Pt(18), BLUE),
        ('', False, Pt(4), DARK),
        ('Step 1：填写基本信息（姓名、角色、科目）', False, Pt(17), DARK),
        ('Step 2：设置授课风格（教学风格、专业领域、引导问题）', False, Pt(17), DARK),
        ('Step 3：上传数字人头像（照片+音频波纹动画）', False, Pt(17), DARK),
        ('Step 4：初始化知识库（一键导入示例资料）', False, Pt(17), DARK),
        ('', False, Pt(4), DARK),
        ('管理面板：助教档案编辑 | 知识库管理 | 学生记忆查看', True, Pt(16), BROWN),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(5500000), Emu(5200000), items)
    add_screenshot_placeholder(slide, Emu(6500000), Emu(1200000), Emu(5200000), Emu(5000000),
                               '【请插入：教师设置向导/管理面板界面截图】')
    print("  Slide 10: 教师管理与数字人")

    # ========================================
    # Slide 11: 应用情况
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '应用情况')

    case_items = [
        ('应用案例', True, Pt(22), BLUE),
        ('成人高校物理教师李老师，在管理员创建账号后，登录系统5分钟完成4步设置：填写姓名和角色、设置专业领域"力学、电磁学"和教学风格"善于用生活例子解释、讲解节奏偏慢"、上传个人照片、一键导入教学资料。', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('在工厂上班的王同学只能在下班后学习，他登录系统选择"李老师"助教，点击电话按钮即可语音对话。AI助教根据记忆主动回顾上次学习进度，直接说话提问即可。对于基础薄弱的王同学，AI助教放慢讲解节奏、用生活化类比解释概念，并在学生表示"懂了"时出题验证。如果王同学问与物理无关的问题，AI助教会礼貌引导回学科内容。', False, Pt(16), DARK),
        ('', False, Pt(6), DARK),
        ('应用成效', True, Pt(22), BLUE),
        ('• 突破时空限制：成人学生碎片时间获得语音辅导，缓解工学矛盾', False, Pt(16), DARK),
        ('• 个性化教学：记住每位学生知识掌握情况和学习偏好，因材施教', False, Pt(16), DARK),
        ('• 价值观约束：AI助教严格聚焦专业领域，守护教育安全底线', False, Pt(16), DARK),
        ('• 教师知识数字化沉淀：教学经验和风格被AI助教永久保留', False, Pt(16), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(7000000), Emu(5500000), case_items)
    add_screenshot_placeholder(slide, Emu(8000000), Emu(1200000), Emu(3800000), Emu(4500000),
                               '【请插入：学生使用场景截图\n（登录选助教/语音对话）】')
    print("  Slide 11: 应用情况")

    # ========================================
    # Slide 12: 创新与展望
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, '创新点与展望')

    innovation_items = [
        ('核心创新', True, Pt(22), BLUE),
        ('• 价值观导向的领域约束：区别于通用AI"有问必答"，坚守教育者角色', False, Pt(18), DARK),
        ('• 回声防护方案：浏览器环境下gain=0静音策略，零硬件解决语音回声', False, Pt(18), DARK),
        ('• 学生记忆驱动因材施教：四维追踪+自动注入，实现"越教越懂学生"', False, Pt(18), DARK),
        ('• 教师知识库深度融入：保留个人教学特色，而非使用通用机器人', False, Pt(18), DARK),
        ('', False, Pt(10), DARK),
        ('未来展望', True, Pt(22), BLUE),
        ('• 引入深度学习VAD模型，提升嘈杂环境鲁棒性（如通勤场景）', False, Pt(18), DARK),
        ('• 增加手写输入、公式识别、拍照提问等多模态交互', False, Pt(18), DARK),
        ('• 加强数据加密与未成年人信息保护合规性', False, Pt(18), DARK),
    ]
    add_textbox(slide, Emu(600000), Emu(1100000), Emu(11000000), Emu(5500000), innovation_items)
    print("  Slide 12: 创新与展望")

    # ========================================
    # Slide 13: 感谢
    # ========================================
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 12192000, 6858000)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x2C, 0x48, 0x9B)
    bg.line.fill.background()
    add_textbox(slide, Emu(800000), Emu(2500000), Emu(10600000), Emu(2000000),
                [('谢谢！', True, Pt(60), WHITE),
                 ('欢迎体验开放智慧助教', False, Pt(24), RGBColor(0xCC, 0xCC, 0xEE))])
    print("  Slide 13: 感谢页")

    # Save
    prs.save(OUTPUT)
    print(f"\nSaved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
